import io
from pptx_engine import replace_text_in_pptx, extract_text_from_pptx
from pptx import Presentation

# create a dummy pptx
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
tf = slide.shapes.title.text_frame
tf.text = "Hello World! This is a test."

buf = io.BytesIO()
prs.save(buf)
buf.seek(0)

# extract text
paragraphs = extract_text_from_pptx(buf)
print("Extracted:", paragraphs)

buf.seek(0)
# replace text
replacements = { "Hello World! This is a test.": "Replaced! It works!" }
rmade, smodified, out_buf = replace_text_in_pptx(buf, replacements)
print("Replacements made:", rmade)
print("Slides modified:", smodified)
