import os
import io
from pptx import Presentation
from typing import List, Tuple, Union

def get_text_frames_from_shape(shape) -> list:
    """
    Recursively extracts all text frames from a shape (including tables and group shapes).
    """
    text_frames = []
    
    if shape.has_text_frame:
        text_frames.append(shape.text_frame)
        
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    text_frames.append(cell.text_frame)
                    
    # If it's a group shape, traverse child shapes
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            text_frames.extend(get_text_frames_from_shape(sub_shape))
            
    return text_frames

def normalize_text(text: str) -> str:
    """
    Normalizes casing, smart quotes, and collapsing extra whitespace 
    to make text matching extremely robust against minor LLM variations.
    """
    if not text:
        return ""
    normalized = text.lower().strip()
    # Normalize quotes
    normalized = normalized.replace("“", '"').replace("”", '"').replace("‘", "'").replace("’", "'")
    # Collapse multiple whitespaces
    normalized = " ".join(normalized.split())
    return normalized

def replace_text_in_pptx(input_file: Union[str, io.BytesIO], replacements: dict) -> Tuple[int, int, io.BytesIO]:
    """
    Ingests a .pptx file or buffer, iterates through all shapes recursively, and replaces text.
    First tries run-level replacement to preserve local formatting and hyperlinks. 
    If the text is split across runs, falls back to merging runs at paragraph level.
    If exact matches fail, falls back to normalized paragraph comparison.
    Returns a tuple of (total_replacements_made, total_slides_modified, output_buffer).
    """
    prs = Presentation(input_file)
    replacements_made = 0
    slides_modified = 0

    for slide in prs.slides:
        slide_had_replacement = False
        text_frames = []
        for shape in slide.shapes:
            text_frames.extend(get_text_frames_from_shape(shape))
            
        for text_frame in text_frames:
            for paragraph in text_frame.paragraphs:
                for old_text, new_text in replacements.items():
                    if not old_text.strip():
                        continue
                    
                    # 1. Try exact matching
                    if old_text in paragraph.text:
                        replaced_at_run_level = False
                        for run in paragraph.runs:
                            if old_text in run.text:
                                # Preserve hyperlinks
                                hlink = run.hyperlink.address if run.hyperlink and run.hyperlink.address else None
                                run.text = run.text.replace(old_text, new_text)
                                if hlink:
                                    run.hyperlink.address = hlink
                                replaced_at_run_level = True
                                replacements_made += 1
                                slide_had_replacement = True
                                print(f"[Exact Match - Run] Replaced: '{old_text}' -> '{new_text}'")
                                break
                        
                        if not replaced_at_run_level:
                            # Merge runs: If it spans multiple runs, merge into the first run to preserve base formatting
                            hlink = paragraph.runs[0].hyperlink.address if paragraph.runs and paragraph.runs[0].hyperlink and paragraph.runs[0].hyperlink.address else None
                            full_text = paragraph.text
                            # Remove all runs except the first
                            for i in range(len(paragraph.runs) - 1, 0, -1):
                                p = paragraph._p
                                p.remove(paragraph.runs[i]._r)
                                
                            if paragraph.runs:
                                run = paragraph.runs[0]
                                run.text = full_text.replace(old_text, new_text)
                                if hlink:
                                    run.hyperlink.address = hlink
                            else:
                                paragraph.text = full_text.replace(old_text, new_text)
                                
                            replacements_made += 1
                            slide_had_replacement = True
                            print(f"[Exact Match - Merged Runs] Replaced: '{old_text}' -> '{new_text}'")
                            
                    # 2. Try normalized paragraph comparison fallback (full match)
                    elif normalize_text(old_text) == normalize_text(paragraph.text):
                        hlink = paragraph.runs[0].hyperlink.address if paragraph.runs and paragraph.runs[0].hyperlink and paragraph.runs[0].hyperlink.address else None
                        
                        # Remove extra runs
                        for i in range(len(paragraph.runs) - 1, 0, -1):
                            p = paragraph._p
                            p.remove(paragraph.runs[i]._r)
                            
                        if paragraph.runs:
                            run = paragraph.runs[0]
                            run.text = new_text
                            if hlink:
                                run.hyperlink.address = hlink
                        else:
                            paragraph.text = new_text
                            
                        replacements_made += 1
                        slide_had_replacement = True
                        print(f"[Normalized Fallback] Replaced paragraph: '{paragraph.text}' -> '{new_text}'")

        if slide_had_replacement:
            slides_modified += 1

    output_buffer = io.BytesIO()
    prs.save(output_buffer)
    output_buffer.seek(0)
    return replacements_made, slides_modified, output_buffer

def validate_pptx(file_buffer: Union[str, io.BytesIO]) -> bool:
    """
    Anti-Corruption Validation: Attempts to re-parse the saved .pptx buffer.
    Returns True if valid, False if corrupted.
    """
    try:
        if isinstance(file_buffer, io.BytesIO):
            file_buffer.seek(0)
        Presentation(file_buffer)
        if isinstance(file_buffer, io.BytesIO):
            file_buffer.seek(0)
        return True
    except Exception as e:
        print(f"Anti-Corruption check failed: {e}")
        return False

def extract_text_from_pptx(input_file: Union[str, io.BytesIO]) -> List[str]:
    """
    Extracts all unique text paragraphs from the presentation recursively to feed into the AI.
    This preserves full sentences, table contents, and grouped text boxes.
    """
    prs = Presentation(input_file)
    paragraphs = []
    
    for slide in prs.slides:
        text_frames = []
        for shape in slide.shapes:
            text_frames.extend(get_text_frames_from_shape(shape))
            
        for text_frame in text_frames:
            for paragraph in text_frame.paragraphs:
                p_text = paragraph.text.strip()
                # Filter out empty text, slide numbers, or extremely short formatting artifacts
                if p_text and len(p_text) > 3 and p_text not in paragraphs:
                    paragraphs.append(p_text)
            
    return paragraphs
