import os
import io
import hashlib
import time
import threading
import sys
from pptx import Presentation
from typing import List, Tuple, Union, Dict

def safe_print(msg: str):
    try:
        print(msg)
    except UnicodeEncodeError:
        try:
            # Fallback to sys.stdout.encoding if available
            encoding = sys.stdout.encoding or "utf-8"
            print(msg.encode(encoding, errors="replace").decode(encoding))
        except Exception:
            try:
                print(msg.encode("ascii", errors="replace").decode("ascii"))
            except Exception:
                pass

# In-memory cache for extracted paragraphs: MD5-hash -> (List[str], timestamp)
_extraction_cache: Dict[str, Tuple[List[str], float]] = {}
EXTRACTION_CACHE_TTL = 300

# In-memory cache for downloaded deck bytes: URL -> (content_bytes, timestamp)
_download_cache: Dict[str, Tuple[bytes, float]] = {}
DOWNLOAD_CACHE_TTL = 300  # 5 minutes cache TTL

# Lock mechanisms for thread-safe caching operations
_download_cache_lock = threading.Lock()
_extraction_cache_lock = threading.Lock()

def get_buffer_hash(buffer: io.BytesIO) -> str:
    """
    Generates MD5 hash for a BytesIO buffer.
    """
    pos = buffer.tell()
    buffer.seek(0)
    content = buffer.read()
    buffer.seek(pos)
    return hashlib.md5(content).hexdigest()

def get_file_cache_key(input_file: Union[str, io.BytesIO]) -> str:
    """
    Determines cache key for the file to prevent redundant parsing.
    """
    if isinstance(input_file, str):
        try:
            mtime = os.path.getmtime(input_file)
            return f"path:{input_file}:{mtime}"
        except Exception:
            return f"path:{input_file}"
    elif isinstance(input_file, io.BytesIO):
        return get_buffer_hash(input_file)
    return None

def get_cached_download(url: str) -> bytes:
    """
    Retrieves downloaded bytes if within TTL, and cleans expired items.
    """
    now = time.time()
    with _download_cache_lock:
        # Prune expired items
        expired = [k for k, v in list(_download_cache.items()) if now - v[1] > DOWNLOAD_CACHE_TTL]
        for k in expired:
            _download_cache.pop(k, None)
            
        if url in _download_cache:
            content, timestamp = _download_cache[url]
            return content
    return None

def set_cached_download(url: str, content: bytes):
    """
    Saves downloaded bytes to memory cache, limiting size to 50 items.
    """
    with _download_cache_lock:
        if len(_download_cache) >= 50:
            # Evict oldest item
            oldest = min(_download_cache.keys(), key=lambda k: _download_cache[k][1])
            _download_cache.pop(oldest, None)
        _download_cache[url] = (content, time.time())

def get_text_frames_from_shape(shape) -> list:
    """
    Recursively extracts all text frames from a shape (including tables and group shapes).
    """
    text_frames = []
    
    if shape.has_text_frame:
        text_frames.append(shape.text_frame)
        
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    text_frames.append(cell.text_frame)
                    
    # If it's a group shape, traverse child shapes
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            text_frames.extend(get_text_frames_from_shape(sub_shape))
            
    return text_frames

def normalize_text(text: str) -> str:
    """
    Normalizes casing, smart quotes, and collapsing extra whitespace 
    to make text matching extremely robust against minor LLM variations.
    """
    if not text:
        return ""
    normalized = text.lower().strip()
    # Normalize quotes
    normalized = normalized.replace("“", '"').replace("”", '"').replace("‘", "'").replace("’", "'")
    # Collapse multiple whitespaces
    normalized = " ".join(normalized.split())
    return normalized

def replace_text_in_pptx(input_file: Union[str, io.BytesIO], replacements: dict) -> Tuple[int, int, io.BytesIO]:
    """
    Ingests a .pptx file or buffer, iterates through all shapes recursively, and replaces text.
    First tries run-level replacement to preserve local formatting and hyperlinks. 
    If the text is split across runs, falls back to merging runs at paragraph level.
    If exact matches fail, falls back to normalized paragraph comparison.
    Returns a tuple of (total_replacements_made, total_slides_modified, output_buffer).
    """
    prs = Presentation(input_file)
    replacements_made = 0
    slides_modified = 0

    # Pre-normalize replacement keys and values to avoid redundant computations in nested loops
    # Store: {old_text: (new_text, normalized_old_text)}
    normalized_replacements = {}
    for old_text, new_text in replacements.items():
        if old_text.strip():
            normalized_replacements[old_text] = (new_text, normalize_text(old_text))

    for slide in prs.slides:
        slide_had_replacement = False
        text_frames = []
        for shape in slide.shapes:
            text_frames.extend(get_text_frames_from_shape(shape))
            
        for text_frame in text_frames:
            for paragraph in text_frame.paragraphs:
                # Cache paragraph.text to avoid multiple expensive property accesses (XML traversal)
                p_text = paragraph.text
                p_text_normalized = None  # Lazy-loaded on demand
                
                for old_text, (new_text, old_text_norm) in normalized_replacements.items():
                    # Check exact match on cached paragraph text
                    if old_text in p_text:
                        replaced_at_run_level = False
                        
                        # Cache run list to prevent re-instantiation
                        runs = list(paragraph.runs)
                        for run in runs:
                            run_text = run.text
                            if old_text in run_text:
                                # Preserve hyperlinks
                                hlink = run.hyperlink.address if run.hyperlink and run.hyperlink.address else None
                                run.text = run_text.replace(old_text, new_text)
                                if hlink:
                                    run.hyperlink.address = hlink
                                replaced_at_run_level = True
                                replacements_made += 1
                                slide_had_replacement = True
                                safe_print(f"[Exact Match - Run] Replaced: '{old_text}' -> '{new_text}'")
                                break
                        
                        if not replaced_at_run_level:
                            # Merge runs: If it spans multiple runs, merge into the first run to preserve base formatting
                            runs = list(paragraph.runs)
                            hlink = runs[0].hyperlink.address if runs and runs[0].hyperlink and runs[0].hyperlink.address else None
                            
                            # Remove all runs except the first
                            for i in range(len(runs) - 1, 0, -1):
                                p = paragraph._p
                                p.remove(runs[i]._r)
                                
                            # Retrieve remaining runs
                            remaining_runs = list(paragraph.runs)
                            if remaining_runs:
                                run = remaining_runs[0]
                                run.text = p_text.replace(old_text, new_text)
                                if hlink:
                                    run.hyperlink.address = hlink
                            else:
                                paragraph.text = p_text.replace(old_text, new_text)
                                
                            replacements_made += 1
                            slide_had_replacement = True
                            safe_print(f"[Exact Match - Merged Runs] Replaced: '{old_text}' -> '{new_text}'")
                        
                        # Update cached paragraph text and invalidate normalized cache
                        p_text = p_text.replace(old_text, new_text)
                        p_text_normalized = None

                    # 2. Try normalized paragraph comparison fallback (full match)
                    else:
                        if p_text_normalized is None:
                            p_text_normalized = normalize_text(p_text)
                            
                        if old_text_norm == p_text_normalized:
                            runs = list(paragraph.runs)
                            hlink = runs[0].hyperlink.address if runs and runs[0].hyperlink and runs[0].hyperlink.address else None
                            
                            # Remove extra runs
                            for i in range(len(runs) - 1, 0, -1):
                                p = paragraph._p
                                p.remove(runs[i]._r)
                                
                            remaining_runs = list(paragraph.runs)
                            if remaining_runs:
                                run = remaining_runs[0]
                                run.text = new_text
                                if hlink:
                                    run.hyperlink.address = hlink
                            else:
                                paragraph.text = new_text
                                
                            replacements_made += 1
                            slide_had_replacement = True
                            safe_print(f"[Normalized Fallback] Replaced paragraph: '{p_text}' -> '{new_text}'")
                            
                            # Update cached paragraph text and invalidate normalized cache
                            p_text = new_text
                            p_text_normalized = None

        if slide_had_replacement:
            slides_modified += 1

    output_buffer = io.BytesIO()
    prs.save(output_buffer)
    output_buffer.seek(0)
    return replacements_made, slides_modified, output_buffer

def validate_pptx(file_buffer: Union[str, io.BytesIO]) -> bool:
    """
    Anti-Corruption Validation: Attempts to re-parse the saved .pptx buffer.
    Returns True if valid, False if corrupted.
    """
    try:
        if isinstance(file_buffer, io.BytesIO):
            file_buffer.seek(0)
        Presentation(file_buffer)
        if isinstance(file_buffer, io.BytesIO):
            file_buffer.seek(0)
        return True
    except Exception as e:
        safe_print(f"Anti-Corruption check failed: {e}")
        return False

def extract_text_from_pptx(input_file: Union[str, io.BytesIO]) -> List[str]:
    """
    Extracts all unique text paragraphs from the presentation recursively to feed into the AI.
    This preserves full sentences, table contents, and grouped text boxes.
    Utilizes an in-memory cache and O(1) set membership for performance.
    """
    # 1. Try to read from cache first
    cache_key = get_file_cache_key(input_file)
    now = time.time()
    if cache_key:
        with _extraction_cache_lock:
            # Prune expired items
            expired = [k for k, v in list(_extraction_cache.items()) if now - v[1] > EXTRACTION_CACHE_TTL]
            for k in expired:
                _extraction_cache.pop(k, None)
                
            if cache_key in _extraction_cache:
                safe_print(f"[Cache Hit] Returning cached paragraphs for {cache_key[:32]}...")
                return _extraction_cache[cache_key][0].copy()

    prs = Presentation(input_file)
    paragraphs = []
    seen = set()
    
    for slide in prs.slides:
        text_frames = []
        for shape in slide.shapes:
            text_frames.extend(get_text_frames_from_shape(shape))
            
        for text_frame in text_frames:
            for paragraph in text_frame.paragraphs:
                p_text = paragraph.text.strip()
                # Filter out empty text, slide numbers, or extremely short formatting artifacts
                if p_text and len(p_text) > 3 and p_text not in seen:
                    seen.add(p_text)
                    paragraphs.append(p_text)
            
    # Cache the result before returning
    if cache_key:
        with _extraction_cache_lock:
            if len(_extraction_cache) >= 50:
                # Evict oldest item
                oldest = min(_extraction_cache.keys(), key=lambda k: _extraction_cache[k][1])
                _extraction_cache.pop(oldest, None)
            _extraction_cache[cache_key] = (paragraphs, time.time())
        
    return paragraphs.copy() if paragraphs else []
